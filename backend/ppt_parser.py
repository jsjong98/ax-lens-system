"""
ppt_parser.py — PPT 파일에서 워크플로우 노드/엣지를 추출

PPT 슬라이드 내의 도형(사각형, 둥근사각형 등)을 노드로,
연결선(커넥터)을 엣지로 파싱합니다.

다양한 PPT 스타일에 대응:
  - 도형 크기/색상 기반 노이즈 필터링
  - 테이블, 플레이스홀더, 장식용 도형 제외
  - 그룹 도형 내부까지 탐색
  - 커넥터 없을 때 위치 기반 순서 추론
"""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ── 필터링 상수 ───────────────────────────────────────────────────────────────

# 노드로 인식할 최소/최대 크기 (인치)
MIN_NODE_WIDTH = 0.5
MIN_NODE_HEIGHT = 0.3
MAX_NODE_WIDTH = 10.0
MAX_NODE_HEIGHT = 5.0

# 노이즈 텍스트 패턴 (제목, 페이지 번호, 날짜 등)
NOISE_PATTERNS = [
    r'^\d{1,3}$',                    # 순수 숫자 (페이지 번호)
    r'^\d{4}[.\-/]\d{1,2}',         # 날짜
    r'^(confidential|draft|내부용)',  # 워터마크
    r'^(page|slide|p\.)\s*\d',       # 페이지 표시
    r'^copyright',                    # 저작권
    r'^\*+$',                        # 별표만
    r'^[\-=_]+$',                    # 구분선 텍스트
]
NOISE_RE = [re.compile(p, re.IGNORECASE) for p in NOISE_PATTERNS]

# 프로세스 노드로 판별할 키워드 (가중치)
PROCESS_KEYWORDS = [
    "수행", "처리", "확인", "검토", "승인", "작성", "분석", "조사",
    "생성", "관리", "설정", "배포", "전달", "수집", "평가", "보고",
    "등록", "실행", "계획", "통보", "접수", "배정", "조율", "협의",
    "산출", "반영", "업데이트", "모니터링", "점검",
]


# ── 데이터 모델 ───────────────────────────────────────────────────────────────

@dataclass
class PptNode:
    """PPT에서 추출한 노드."""
    id: str
    text: str
    slide_index: int
    left: float      # inches
    top: float        # inches
    width: float      # inches
    height: float     # inches
    shape_type: str = ""
    fill_color: str = ""
    font_size: float = 0.0   # pt
    is_title: bool = False
    node_confidence: float = 0.0  # 프로세스 노드일 확률

    @property
    def center_x(self) -> float:
        return self.left + self.width / 2

    @property
    def center_y(self) -> float:
        return self.top + self.height / 2


@dataclass
class PptEdge:
    """PPT에서 추출한 연결선."""
    id: str
    source_id: str | None = None
    target_id: str | None = None
    label: str = ""
    slide_index: int = 0


@dataclass
class PptSlide:
    """파싱된 슬라이드."""
    index: int
    title: str = ""
    nodes: list[PptNode] = field(default_factory=list)
    edges: list[PptEdge] = field(default_factory=list)


@dataclass
class ParsedPpt:
    """PPT 파싱 결과 전체."""
    filename: str
    slide_count: int
    slides: list[PptSlide] = field(default_factory=list)


# ── 파서 ──────────────────────────────────────────────────────────────────────

def parse_ppt(source: str | Path | bytes | BytesIO) -> ParsedPpt:
    """PPT 파일을 파싱하여 슬라이드별 노드/엣지를 추출합니다."""
    if isinstance(source, bytes):
        source = BytesIO(source)
    prs = Presentation(source)

    filename = ""
    if isinstance(source, (str, Path)):
        filename = Path(source).name

    slide_width = _emu_to_inches(prs.slide_width) if prs.slide_width else 13.33
    slide_height = _emu_to_inches(prs.slide_height) if prs.slide_height else 7.5

    slides: list[PptSlide] = []

    for slide_idx, slide in enumerate(prs.slides):
        ppt_slide = PptSlide(index=slide_idx)

        # 슬라이드 제목
        if slide.shapes.title:
            ppt_slide.title = slide.shapes.title.text.strip()

        all_shapes: list[tuple[Any, bool]] = []  # (shape, is_in_group)
        node_map: dict[int, PptNode] = {}

        # 모든 도형을 평탄화 (그룹 포함)
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    all_shapes.append((child, True))
            else:
                all_shapes.append((shape, False))

        # 1차: 커넥터와 노드 분류
        for shape, is_grouped in all_shapes:
            if _is_connector(shape):
                edge = _parse_connector(shape, slide_idx, len(ppt_slide.edges))
                if edge:
                    ppt_slide.edges.append(edge)
                continue

            # 테이블, 차트, 미디어 등 제외
            if _should_skip_shape(shape):
                continue

            node = _parse_shape_as_node(
                shape, slide_idx, len(node_map),
                slide_width, slide_height,
                is_title_shape=(shape == slide.shapes.title if slide.shapes.title else False),
            )
            if node and node.node_confidence > 0.3:
                node_map[shape.shape_id] = node

        ppt_slide.nodes = list(node_map.values())

        # 2차: 커넥터의 source/target을 shape_id → node_id로 매핑
        shape_id_to_node_id = {sid: n.id for sid, n in node_map.items()}
        for edge in ppt_slide.edges:
            if edge.source_id and edge.source_id.isdigit():
                edge.source_id = shape_id_to_node_id.get(int(edge.source_id), edge.source_id)
            if edge.target_id and edge.target_id.isdigit():
                edge.target_id = shape_id_to_node_id.get(int(edge.target_id), edge.target_id)

        # 유효한 엣지만 유지 (양쪽 노드가 모두 존재)
        valid_node_ids = {n.id for n in ppt_slide.nodes}
        ppt_slide.edges = [
            e for e in ppt_slide.edges
            if e.source_id in valid_node_ids and e.target_id in valid_node_ids
        ]

        # 엣지가 없으면 위치 기반 순서 추론
        if not ppt_slide.edges and len(ppt_slide.nodes) > 1:
            ppt_slide.edges = _infer_edges_from_position(ppt_slide.nodes, slide_idx)

        slides.append(ppt_slide)

    return ParsedPpt(
        filename=filename,
        slide_count=len(slides),
        slides=slides,
    )


# ── 도형 판별 ─────────────────────────────────────────────────────────────────

def _should_skip_shape(shape: Any) -> bool:
    """노드 후보에서 제외할 도형 유형."""
    try:
        st = shape.shape_type
        # 테이블, 차트, 미디어, 그림, OLE 등은 제외
        skip_types = {
            MSO_SHAPE_TYPE.TABLE,
            MSO_SHAPE_TYPE.CHART,
            MSO_SHAPE_TYPE.MEDIA,
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
            MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
        }
        if st in skip_types:
            return True

        # 플레이스홀더 중 제목/부제/날짜/슬라이드번호는 제외
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            # 0=TITLE, 1=CENTER_TITLE, 2=SUBTITLE, 10=DATE, 11=SLIDE_NUMBER, 12=FOOTER
            if ph_type in (0, 1, 2, 10, 11, 12):
                return True
    except Exception:
        pass
    return False


def _is_connector(shape: Any) -> bool:
    """도형이 연결선(커넥터)인지 판별합니다."""
    try:
        # python-pptx의 커넥터 타입 직접 체크
        if hasattr(shape, 'begin_x') and hasattr(shape, 'end_x'):
            return True
        # XML 태그 확인 — cxnSp = 커넥터 도형
        tag = shape._element.tag if hasattr(shape, '_element') else ""
        if 'cxnSp' in tag:
            return True
    except Exception:
        pass
    return False


def _parse_connector(shape: Any, slide_idx: int, edge_idx: int) -> PptEdge | None:
    """커넥터 도형에서 엣지 정보를 추출합니다."""
    try:
        edge_id = f"edge-s{slide_idx}-{edge_idx}"
        source_id = None
        target_id = None

        el = shape._element
        for cxn_tag in el.iter():
            tag_name = cxn_tag.tag.split('}')[-1] if '}' in cxn_tag.tag else cxn_tag.tag
            if tag_name == 'stCxn':
                source_id = cxn_tag.get('id')
            elif tag_name == 'endCxn':
                target_id = cxn_tag.get('id')

        label = ""
        if hasattr(shape, 'text') and shape.text:
            label = shape.text.strip()

        return PptEdge(
            id=edge_id,
            source_id=source_id,
            target_id=target_id,
            label=label,
            slide_index=slide_idx,
        )
    except Exception:
        return None


def _parse_shape_as_node(
    shape: Any,
    slide_idx: int,
    node_idx: int,
    slide_width: float,
    slide_height: float,
    is_title_shape: bool = False,
) -> PptNode | None:
    """도형을 노드로 변환합니다. 프로세스 노드 확률도 계산."""
    try:
        if not hasattr(shape, 'text') or not shape.text.strip():
            return None

        text = shape.text.strip()
        # 줄바꿈을 공백으로 정규화
        text_norm = re.sub(r'\s+', ' ', text).strip()

        if len(text_norm) < 2:
            return None

        # 노이즈 패턴 체크
        for pattern in NOISE_RE:
            if pattern.search(text_norm):
                return None

        # 크기
        left = _emu_to_inches(shape.left) if shape.left else 0
        top = _emu_to_inches(shape.top) if shape.top else 0
        width = _emu_to_inches(shape.width) if shape.width else 0
        height = _emu_to_inches(shape.height) if shape.height else 0

        # 크기 필터
        if width < MIN_NODE_WIDTH or height < MIN_NODE_HEIGHT:
            return None
        if width > MAX_NODE_WIDTH or height > MAX_NODE_HEIGHT:
            return None

        # 폰트 크기 추출
        font_size = _extract_font_size(shape)

        # 채우기 색상
        fill_color = _extract_fill_color(shape)

        # 도형 유형
        shape_type = ""
        try:
            shape_type = str(shape.shape_type)
        except Exception:
            pass

        # ── 프로세스 노드 확률 계산 ──
        confidence = _calculate_node_confidence(
            text=text_norm,
            width=width,
            height=height,
            font_size=font_size,
            fill_color=fill_color,
            shape_type=shape_type,
            is_title=is_title_shape,
            slide_width=slide_width,
        )

        node_id = f"ppt-s{slide_idx}-n{node_idx}"

        return PptNode(
            id=node_id,
            text=text_norm,
            slide_index=slide_idx,
            left=left,
            top=top,
            width=width,
            height=height,
            shape_type=shape_type,
            fill_color=fill_color,
            font_size=font_size,
            is_title=is_title_shape,
            node_confidence=confidence,
        )
    except Exception:
        return None


def _calculate_node_confidence(
    text: str,
    width: float,
    height: float,
    font_size: float,
    fill_color: str,
    shape_type: str,
    is_title: bool,
    slide_width: float,
) -> float:
    """도형이 프로세스 노드일 확률을 0~1로 계산합니다."""
    score = 0.5  # 기본값

    # 제목이면 노드 아님
    if is_title:
        return 0.0

    # 너무 큰 도형 (슬라이드 너비의 70% 이상) — 배경/장식일 가능성
    if width > slide_width * 0.7:
        score -= 0.3

    # 적절한 크기 (1~4인치) — 전형적인 프로세스 노드
    if 0.8 <= width <= 4.0 and 0.3 <= height <= 1.5:
        score += 0.2

    # 텍스트가 너무 길면 (100자 이상) 설명 텍스트일 가능성
    if len(text) > 100:
        score -= 0.2
    elif len(text) > 50:
        score -= 0.1

    # 적절한 길이 (3~30자)
    if 3 <= len(text) <= 30:
        score += 0.1

    # 프로세스 키워드 포함
    keyword_count = sum(1 for kw in PROCESS_KEYWORDS if kw in text)
    score += min(keyword_count * 0.1, 0.3)

    # ID 패턴 포함 (1.1.1 같은)
    if re.search(r'\d+\.\d+\.\d+', text):
        score += 0.2

    # 채우기 색상 있으면 노드일 가능성 높음
    if fill_color:
        score += 0.1

    # 폰트 크기가 너무 크면 제목/헤더
    if font_size > 20:
        score -= 0.2
    elif font_size > 14:
        score -= 0.1

    return max(0.0, min(1.0, score))


# ── 유틸리티 ──────────────────────────────────────────────────────────────────

def _emu_to_inches(emu: int | None) -> float:
    if emu is None:
        return 0.0
    return emu / 914400


def _extract_fill_color(shape: Any) -> str:
    try:
        fill = shape.fill
        if fill and fill.type is not None:
            fc = fill.fore_color
            if fc and fc.rgb:
                return str(fc.rgb)
    except Exception:
        pass
    return ""


def _extract_font_size(shape: Any) -> float:
    """도형 내 첫 번째 텍스트의 폰트 크기를 pt로 반환합니다."""
    try:
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font and run.font.size:
                        return run.font.size.pt
    except Exception:
        pass
    return 0.0


def _infer_edges_from_position(
    nodes: list[PptNode], slide_idx: int
) -> list[PptEdge]:
    """
    커넥터가 없을 때 노드 위치 기반으로 순서를 추론합니다.
    Y좌표 기준 레벨 그룹핑 후, 레벨 간 순차 엣지 생성.
    """
    Y_THRESHOLD = 0.5  # inches

    sorted_nodes = sorted(nodes, key=lambda n: (n.center_y, n.center_x))

    levels: list[list[PptNode]] = []
    current_level: list[PptNode] = [sorted_nodes[0]]
    current_y = sorted_nodes[0].center_y

    for node in sorted_nodes[1:]:
        if abs(node.center_y - current_y) <= Y_THRESHOLD:
            current_level.append(node)
        else:
            levels.append(current_level)
            current_level = [node]
            current_y = node.center_y
    levels.append(current_level)

    edges: list[PptEdge] = []
    edge_idx = 0

    for i in range(len(levels) - 1):
        src_level = levels[i]
        tgt_level = levels[i + 1]

        if len(src_level) == 1 and len(tgt_level) == 1:
            # 1:1 연결
            edges.append(PptEdge(
                id=f"inferred-s{slide_idx}-{edge_idx}",
                source_id=src_level[0].id,
                target_id=tgt_level[0].id,
                slide_index=slide_idx,
            ))
            edge_idx += 1
        elif len(src_level) == 1:
            # 1:N (분기)
            for tgt in tgt_level:
                edges.append(PptEdge(
                    id=f"inferred-s{slide_idx}-{edge_idx}",
                    source_id=src_level[0].id,
                    target_id=tgt.id,
                    slide_index=slide_idx,
                ))
                edge_idx += 1
        elif len(tgt_level) == 1:
            # N:1 (합류)
            for src in src_level:
                edges.append(PptEdge(
                    id=f"inferred-s{slide_idx}-{edge_idx}",
                    source_id=src.id,
                    target_id=tgt_level[0].id,
                    slide_index=slide_idx,
                ))
                edge_idx += 1
        else:
            # N:M — X좌표 가까운 것끼리 매칭
            for src in src_level:
                closest_tgt = min(tgt_level, key=lambda t: abs(t.center_x - src.center_x))
                edges.append(PptEdge(
                    id=f"inferred-s{slide_idx}-{edge_idx}",
                    source_id=src.id,
                    target_id=closest_tgt.id,
                    slide_index=slide_idx,
                ))
                edge_idx += 1

    return edges


# ── 노드-태스크 매칭 ──────────────────────────────────────────────────────────

def match_nodes_to_tasks(
    nodes: list[PptNode],
    tasks: list[dict],
) -> list[dict]:
    """PPT 노드 텍스트를 태스크 목록과 매칭합니다."""
    results = []
    for node in nodes:
        best_match = _find_best_match(node.text, tasks)
        results.append({
            "node_id": node.id,
            "node_text": node.text,
            "position": {"x": node.center_x, "y": node.center_y},
            "node_confidence": round(node.node_confidence, 2),
            "matched_task_id": best_match["id"] if best_match else None,
            "matched_task_name": best_match["name"] if best_match else None,
            "matched_level": best_match.get("level", "") if best_match else None,
            "match_confidence": best_match["score"] if best_match else 0,
        })

    return results


def _find_best_match(node_text: str, tasks: list[dict]) -> dict | None:
    """
    노드 텍스트에 가장 잘 매칭되는 태스크를 찾습니다.

    매칭 전략 (이름 우선):
      1. 이름 완전 일치 / 포함 관계 → 높은 점수
      2. 단어(키워드) 겹침 비율 → 중간 점수
      3. 동점이면 ID가 텍스트에 포함된 쪽 우선 (중복 구분용)
    """
    text = re.sub(r'\s+', ' ', node_text).strip()
    if not text:
        return None

    norm_text = text.lower()
    # 한글/영어 단어 추출 (2자 이상)
    text_words = set(w for w in re.findall(r'[가-힣a-zA-Z]{2,}', norm_text))

    best: dict | None = None
    best_score = 0.0

    for task in tasks:
        task_id = task.get("id", "")
        task_name = task.get("name", "")
        task_l4 = task.get("l4", "")
        task_l4_id = task.get("l4_id", "")

        score = 0.0
        id_bonus = 0.0

        # ── 이름 기반 매칭 (L5 이름 우선, L4 이름 보조) ──
        norm_name = re.sub(r'\s+', ' ', task_name).strip().lower() if task_name else ""
        norm_l4 = re.sub(r'\s+', ' ', task_l4).strip().lower() if task_l4 else ""

        # 1-a. L5 이름 완전 일치
        if norm_name and norm_name == norm_text:
            score = 0.95
        # 1-b. L4 이름 완전 일치
        elif norm_l4 and norm_l4 == norm_text:
            score = 0.90
        # 2-a. L5 이름 포함 관계 (한쪽이 다른 쪽에 포함)
        elif norm_name and (norm_name in norm_text or norm_text in norm_name):
            shorter = min(len(norm_name), len(norm_text))
            longer = max(len(norm_name), len(norm_text))
            score = 0.5 + (shorter / max(longer, 1)) * 0.40
        # 2-b. L4 이름 포함 관계
        elif norm_l4 and (norm_l4 in norm_text or norm_text in norm_l4):
            shorter = min(len(norm_l4), len(norm_text))
            longer = max(len(norm_l4), len(norm_text))
            score = 0.4 + (shorter / max(longer, 1)) * 0.35
        else:
            # 3. 단어 겹침 비율 (키워드 매칭)
            name_words = set(w for w in re.findall(r'[가-힣a-zA-Z]{2,}', norm_name))
            if name_words and text_words:
                overlap = len(name_words & text_words)
                total = len(name_words | text_words)
                jaccard = overlap / max(total, 1)
                if jaccard > 0:
                    score = 0.3 + jaccard * 0.50

        # ── ID 보너스 (동점 구분용, 이름 매칭이 있을 때만) ──
        if score > 0:
            if task_id and task_id in text:
                id_bonus = 0.05
            elif task_l4_id and task_l4_id in text:
                id_bonus = 0.03

        final_score = min(score + id_bonus, 1.0)

        if final_score > best_score:
            best_score = final_score
            level = "L5" if task_id.count(".") >= 3 else "L4"
            best = {"id": task_id, "name": task_name, "score": round(final_score, 2), "level": level}

    # 임계값: 0.3 이상이면 매칭 (이전 0.4보다 관대)
    if best and best_score < 0.3:
        return None
    return best


# ── React Flow 변환 ───────────────────────────────────────────────────────────

def ppt_slide_to_react_flow(slide: PptSlide, scale: float = 100.0) -> dict:
    """PptSlide를 React Flow 호환 JSON으로 변환합니다."""
    nodes = []
    for node in slide.nodes:
        nodes.append({
            "id": node.id,
            "type": "l4",
            "position": {"x": node.left * scale, "y": node.top * scale},
            "data": {
                "label": node.text,
                "level": "L4",
                "id": "",
                "description": "",
                "source": "ppt",
                "nodeConfidence": node.node_confidence,
            },
        })

    edges = []
    for edge in slide.edges:
        if edge.source_id and edge.target_id:
            edges.append({
                "id": edge.id,
                "source": edge.source_id,
                "target": edge.target_id,
                "type": "smoothstep",
                "animated": True,
                "label": edge.label,
                "style": {"stroke": "#a62121", "strokeWidth": 2},
                "markerEnd": {"type": "arrowclosed", "width": 20, "height": 20, "color": "#a62121"},
            })

    return {"nodes": nodes, "edges": edges, "slide_title": slide.title}


# ── PPT → WorkflowSheet 변환 ──────────────────────────────────────────────

def ppt_to_parsed_workflow(
    parsed_ppt: ParsedPpt,
    matches_per_slide: list[list[dict]] | None = None,
) -> "ParsedWorkflow":
    """
    PPT 파싱 결과를 workflow_parser의 ParsedWorkflow 형태로 변환합니다.
    이를 통해 To-Be 생성 등 JSON 워크플로우와 동일한 파이프라인을 사용할 수 있습니다.

    matches_per_slide: 슬라이드별 match_nodes_to_tasks() 결과 (있으면 task_id/level 정보 사용)
    """
    from workflow_parser import (
        WorkflowNode, WorkflowEdge, ExecutionStep,
        WorkflowSheet, ParsedWorkflow,
    )

    sheets: list[WorkflowSheet] = []

    for si, slide in enumerate(parsed_ppt.slides):
        if not slide.nodes:
            continue

        matches = matches_per_slide[si] if matches_per_slide and si < len(matches_per_slide) else []
        match_map = {m["node_id"]: m for m in matches} if matches else {}

        # 노드 변환
        wf_nodes: dict[str, WorkflowNode] = {}
        for node in slide.nodes:
            m = match_map.get(node.id, {})
            matched_id = m.get("matched_task_id", "")
            matched_level = m.get("matched_level", "")

            # 매칭된 경우 task_id와 level 사용, 아니면 PPT 텍스트 기반 추정
            task_id = matched_id or node.id
            level = matched_level or _guess_level(matched_id)

            wf_nodes[node.id] = WorkflowNode(
                id=node.id,
                level=level,
                task_id=task_id,
                label=node.text,
                description=m.get("matched_task_name", ""),
                position_x=node.center_x,
                position_y=node.center_y,
            )

        # 엣지 변환
        wf_edges: list[WorkflowEdge] = []
        for edge in slide.edges:
            if edge.source_id and edge.target_id:
                wf_edges.append(WorkflowEdge(
                    id=edge.id,
                    source=edge.source_id,
                    target=edge.target_id,
                    label=edge.label,
                ))

        # 실행 순서 분석
        execution_order = _analyze_ppt_execution_order(slide.nodes, slide.edges)

        sheet = WorkflowSheet(
            sheet_id=f"ppt-slide-{si}",
            name=slide.title or f"슬라이드 {si + 1}",
            nodes=wf_nodes,
            edges=wf_edges,
            execution_order=execution_order,
        )
        sheets.append(sheet)

    return ParsedWorkflow(version="ppt-1.0", sheets=sheets)


def _guess_level(task_id: str) -> str:
    """task_id의 점(.) 개수로 레벨을 추정합니다."""
    if not task_id:
        return "L4"
    dots = task_id.count(".")
    if dots >= 4:
        return "L5"
    elif dots >= 3:
        return "L5"
    elif dots >= 2:
        return "L4"
    elif dots >= 1:
        return "L3"
    return "L4"


def _analyze_ppt_execution_order(
    nodes: list[PptNode],
    edges: list[PptEdge],
) -> list["ExecutionStep"]:
    """PPT 노드/엣지에서 실행 순서를 분석합니다."""
    from workflow_parser import ExecutionStep
    from collections import defaultdict, deque

    if not nodes:
        return []

    node_ids = {n.id for n in nodes}

    # 유효한 엣지만 사용
    valid_edges = [e for e in edges if e.source_id in node_ids and e.target_id in node_ids]

    if valid_edges:
        # Kahn's algorithm (topological sort)
        in_degree: dict[str, int] = {nid: 0 for nid in node_ids}
        successors: dict[str, list[str]] = defaultdict(list)
        for e in valid_edges:
            in_degree[e.target_id] = in_degree.get(e.target_id, 0) + 1
            successors[e.source_id].append(e.target_id)

        queue = deque(nid for nid, deg in in_degree.items() if deg == 0)
        steps: list[ExecutionStep] = []
        step_num = 1

        while queue:
            level_nodes = sorted(queue)
            steps.append(ExecutionStep(
                step_number=step_num,
                node_ids=level_nodes,
                is_parallel=len(level_nodes) > 1,
            ))
            step_num += 1

            next_queue: list[str] = []
            for nid in level_nodes:
                for succ in successors.get(nid, []):
                    in_degree[succ] -= 1
                    if in_degree[succ] == 0:
                        next_queue.append(succ)
            queue = deque(next_queue)

        return steps
    else:
        # 엣지 없으면 Y좌표 기반 정렬
        sorted_nodes = sorted(nodes, key=lambda n: (n.center_y, n.center_x))
        steps: list[ExecutionStep] = []
        for i, node in enumerate(sorted_nodes):
            steps.append(ExecutionStep(
                step_number=i + 1,
                node_ids=[node.id],
                is_parallel=False,
            ))
        return steps
