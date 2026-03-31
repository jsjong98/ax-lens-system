"""
ppt_exporter.py — 과제 정의서/설계서 PPT 내보내기

템플릿 PPT를 열어 데이터를 채운 뒤 BytesIO로 반환합니다.
템플릿 구조:
  Slide 1: 표지
  Slide 2: 과제 정의서 (섹션 1~6)
  Slide 3: 과제 설계서 (섹션 7~9)
  Slide 4: 별첨: Agent 정의서 (Agent별 복제)
  Slide 5: Thank You
"""
from __future__ import annotations

import copy
import io
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree


TEMPLATE_PATH = Path(__file__).parent / "template.pptx"

PWC_RED = RGBColor(0xA6, 0x21, 0x21)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x55, 0x55, 0x55)
DARK = RGBColor(0x33, 0x33, 0x33)


# ── 유틸 ──────────────────────────────────────────────────────────────────────

def _find_shape(slide_or_group, name: str):
    """슬라이드/그룹에서 이름으로 shape 검색 (재귀)."""
    shapes = slide_or_group.shapes if hasattr(slide_or_group, 'shapes') else []
    for shape in shapes:
        if shape.name == name:
            return shape
        if hasattr(shape, 'shapes'):
            found = _find_shape(shape, name)
            if found:
                return found
    return None


def _find_shapes_containing(slide_or_group, text_fragment: str) -> list:
    """텍스트를 포함하는 모든 shape 검색 (재귀)."""
    results = []
    shapes = slide_or_group.shapes if hasattr(slide_or_group, 'shapes') else []
    for shape in shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text
            if text_fragment in full_text:
                results.append(shape)
        if hasattr(shape, 'shapes'):
            results.extend(_find_shapes_containing(shape, text_fragment))
    return results


def _set_text(shape, text: str, font_size: Pt | None = None, bold: bool | None = None,
              color: RGBColor | None = None, alignment=None):
    """Shape의 텍스트를 교체합니다. 기존 포맷을 최대한 보존합니다."""
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    if tf.paragraphs:
        para = tf.paragraphs[0]
        if para.runs:
            # 첫 run에 텍스트 설정
            run = para.runs[0]
            run.text = text
            if font_size is not None:
                run.font.size = font_size
            if bold is not None:
                run.font.bold = bold
            if color is not None:
                run.font.color.rgb = color
            # 나머지 run 제거 (중복 텍스트 방지)
            for extra_run in list(para.runs)[1:]:
                extra_run._r.getparent().remove(extra_run._r)
        else:
            para.text = text
        # 나머지 paragraph 제거
        for extra_para in list(tf.paragraphs)[1:]:
            extra_para._p.getparent().remove(extra_para._p)
        if alignment is not None:
            para.alignment = alignment


def _set_multiline_text(shape, lines: list[str], font_size=Pt(11), bold=False,
                        color=DARK, bullet_char="", line_spacing=1.2):
    """Shape에 여러 줄 텍스트를 설정합니다. 줄 수에 따라 폰트 자동 축소."""
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    # 줄 수에 따라 폰트 크기 자동 조절
    total_chars = sum(len(l) for l in lines)
    if len(lines) > 10 or total_chars > 300:
        font_size = Pt(6)
    elif len(lines) > 7 or total_chars > 200:
        font_size = Pt(7)
    elif len(lines) > 4 or total_chars > 120:
        font_size = Pt(8)
    elif len(lines) > 2:
        font_size = Pt(9)

    import re
    from pptx.oxml.ns import qn
    from lxml import etree

    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        # 템플릿 XML의 불릿 설정 제거 (더블 불릿 방지)
        pPr = para._p.get_or_add_pPr()
        for bu_tag in ['a:buChar', 'a:buAutoNum', 'a:buFont', 'a:buSzPct', 'a:buSzPts']:
            for child in pPr.findall(qn(bu_tag)):
                pPr.remove(child)
        etree.SubElement(pPr, qn('a:buNone'))

        # 기존 bullet/기호 제거 후 bullet_char 추가 (중복 방지)
        clean = re.sub(r'^[\s•·‧∙●○◦◉►▶▸▪▫■□◆◇\-–—»›※★☆✓✔·]+', '', line).strip()
        if bullet_char and clean:
            full_text = f"{bullet_char}{clean}"
        else:
            full_text = clean
        run = para.add_run()
        run.text = full_text
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color

        if line_spacing != 1.0:
            para.line_spacing = line_spacing


def _add_text_box(slide, left, top, width, height, text: str,
                  font_size=Pt(11), bold=False, color=DARK, word_wrap=True,
                  alignment=PP_ALIGN.LEFT):
    """슬라이드에 텍스트 박스를 추가합니다."""
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    para = tf.paragraphs[0]
    para.alignment = alignment
    run = para.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _add_multiline_textbox(slide, left, top, width, height, lines: list[str],
                           font_size=Pt(10), color=DARK, bullet="•", line_spacing=1.15,
                           font_name: str = ""):
    """여러 줄 텍스트 박스를 추가합니다."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    import re
    from pptx.oxml.ns import qn
    from lxml import etree

    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        # XML 불릿 설정 제거 (더블 불릿 방지)
        pPr = para._p.get_or_add_pPr()
        for bu_tag in ['a:buChar', 'a:buAutoNum', 'a:buFont', 'a:buSzPct', 'a:buSzPts']:
            for child in pPr.findall(qn(bu_tag)):
                pPr.remove(child)
        etree.SubElement(pPr, qn('a:buNone'))

        # 기존 bullet 제거 후 새로 추가 (중복 방지)
        clean = re.sub(r'^[\s•·‧∙●○◦◉►▶▸▪▫■□◆◇\-–—»›※★☆✓✔·]+', '', line).strip()
        if bullet and clean:
            full_text = f"{bullet}{clean}"
        else:
            full_text = clean
        run = para.add_run()
        run.text = full_text
        run.font.size = font_size
        run.font.color.rgb = color
        if font_name:
            run.font.name = font_name
        para.line_spacing = line_spacing

    return txBox


def _add_grouped_textbox(slide, left, top, width, height, items: list[str],
                          font_size=Pt(7), color=DARK, line_spacing=1.0):
    """방법론용 텍스트 박스 — 각 항목 내 줄바꿈(\n)을 지원하며 들여쓰기 적용."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first_para = True

    for item in items:
        sub_lines = item.split("\n")
        for k, sub in enumerate(sub_lines):
            if first_para:
                para = tf.paragraphs[0]
                first_para = False
            else:
                para = tf.add_paragraph()

            run = para.add_run()
            run.text = sub
            # 들여쓰기된 줄은 더 작은 폰트 + 회색
            if k > 0:
                run.font.size = Pt(max(font_size.pt - 1, 5))
                run.font.color.rgb = RGBColor(0x88, 0x87, 0x80)
                para.line_spacing = 0.95
            else:
                run.font.size = font_size
                run.font.color.rgb = color
                para.line_spacing = line_spacing

    return txBox


def _duplicate_slide(prs: Presentation, slide_index: int) -> Any:
    """슬라이드를 복제합니다 (관계 + XML + 이미지 전체 복사)."""
    from lxml import etree
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    template_slide = prs.slides[slide_index]
    slide_layout = template_slide.slide_layout

    # 새 슬라이드 추가
    new_slide = prs.slides.add_slide(slide_layout)

    # 새 슬라이드의 기존 shape 모두 제거
    spTree = new_slide.shapes._spTree
    for sp in list(spTree):
        tag = sp.tag.split('}')[-1] if '}' in sp.tag else sp.tag
        if tag in ('sp', 'grpSp', 'pic', 'cxnSp', 'graphicFrame'):
            spTree.remove(sp)

    # 원본 슬라이드의 이미지/미디어 관계를 동일한 rId로 복사
    existing_rids = {r.rId for r in new_slide.part.rels.values()}
    for rel in template_slide.part.rels.values():
        if "image" in rel.reltype or "media" in str(rel.target_ref):
            if rel.rId not in existing_rids:
                # python-pptx 내부 API로 동일 rId 유지
                new_slide.part.rels._rels[rel.rId] = rel

    # 원본 슬라이드의 전체 spTree children을 deepcopy
    src_spTree = template_slide.shapes._spTree
    for child in src_spTree:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'grpSp', 'pic', 'cxnSp', 'graphicFrame'):
            spTree.append(copy.deepcopy(child))

    return new_slide


# ── 슬라이드 1: 표지 ──────────────────────────────────────────────────────────

def _fill_cover_slide(slide, project_title: str):
    """표지 슬라이드의 텍스트를 업데이트합니다."""
    # "AI 과제정의서" 텍스트가 있는 shape 찾기
    subtitle = _find_shape(slide, "Subtitle 3")
    if subtitle and subtitle.has_text_frame:
        for para in subtitle.text_frame.paragraphs:
            if "과제정의서" in para.text:
                for run in para.runs:
                    if "과제정의서" in run.text:
                        run.text = "AI 과제정의서"


# ── 슬라이드 2: 과제 정의서 ────────────────────────────────────────────────────

def _fill_definition_slide(slide, definition: dict):
    """과제 정의서 슬라이드에 데이터를 채웁니다."""
    project_number = definition.get("project_number", "")
    project_title = definition.get("project_title", "")
    created_date = definition.get("created_date", "")
    author = definition.get("author", "")

    # 헤더: "■ 과제번호/이름" 뒤에 제목 넣기
    header_shape = _find_shape(slide, "직사각형 7")
    if header_shape:
        _set_text(header_shape, f"■ 과제번호/이름     {project_number}. {project_title}")

    # 그룹 12 안의 입력 필드 텍스트 제거 (기존 "AI 자동화 과제/이름을입력해주세요" 제거)
    group12 = _find_shape(slide, "그룹 12")
    if group12 and hasattr(group12, 'shapes'):
        for sub in group12.shapes:
            if sub.has_text_frame:
                for para in sub.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""

    # 작성일/작성자
    date_shape = _find_shape(slide, "직사각형 13")
    if date_shape:
        _set_text(date_shape, f"작성일: {created_date}   |   작성자: {author}")

    # 1. 과제 개요 — "ㅇㅇ" 플레이스홀더 → overview 리스트
    overview = definition.get("overview", [])
    overview_shape = _find_shape(slide, "직사각형 26")
    if overview_shape and overview:
        _set_multiline_text(overview_shape, overview, font_size=Pt(10), bullet_char="•", color=DARK)

    # 2. 매핑 프로세스 테이블
    table_shape = _find_shape(slide, "표 38")
    if table_shape and table_shape.has_table:
        tbl = table_shape.table
        processes = definition.get("mapping_processes", [])

        # 기존 행 삭제 후 데이터 채우기 (헤더행 + 데이터행)
        # 테이블은 5행 2열, 첫 행이 헤더
        for ri in range(1, min(len(processes) + 1, 5)):
            mp = processes[ri - 1] if ri - 1 < len(processes) else {}
            tbl.cell(ri, 0).text = mp.get("no", "")
            tbl.cell(ri, 1).text = mp.get("process_name", "")
            # task_range가 있으면 프로세스명 아래에 추가
            task_range = mp.get("task_range", "")
            if task_range:
                tbl.cell(ri, 1).text = f"{mp.get('process_name', '')}\n{task_range}"

            # 셀 폰트 설정
            for ci in range(2):
                cell = tbl.cell(ri, ci)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(9)
                        run.font.color.rgb = DARK

    # 3. 이해관계자
    stakeholder = definition.get("stakeholder", {})
    # 과제 오너
    owner_val = _find_shape(slide, "직사각형 48")
    if owner_val:
        _set_text(owner_val, f": {stakeholder.get('project_owner', '')}")
    # 주관 부서
    dept_val = _find_shape(slide, "직사각형 52")
    if dept_val:
        _set_text(dept_val, f": {stakeholder.get('owner_department', '')}")
    # 협업 부서
    collab_val = _find_shape(slide, "직사각형 54")
    if collab_val:
        depts = stakeholder.get("collaborating_departments", [])
        _set_text(collab_val, f": {', '.join(depts) if depts else '없음'}")
    # 외부 파트너
    partner_val = _find_shape(slide, "직사각형 56")
    if partner_val:
        partners = stakeholder.get("external_partners", [])
        _set_text(partner_val, f": {', '.join(partners) if partners else '없음'}")

    # 4. 현황 및 문제점
    cvi = definition.get("current_vs_improvement", {})
    issues_shape = _find_shape(slide, "직사각형 74")
    if issues_shape:
        issues = cvi.get("current_issues", [])
        _set_multiline_text(issues_shape, issues, font_size=Pt(9), bullet_char="•", color=DARK)

    # 4. 개선 방향
    improve_shape = _find_shape(slide, "직사각형 75")
    if improve_shape:
        directions = cvi.get("improvement_directions", [])
        _set_multiline_text(improve_shape, directions, font_size=Pt(9), bullet_char="•", color=DARK)

    # 5. 기대 효과 — 정량적
    quant_shape = _find_shape(slide, "직사각형 99")
    if quant_shape:
        quants = definition.get("expected_effects", {}).get("quantitative", [])
        _set_multiline_text(quant_shape, quants, font_size=Pt(9), bullet_char="•", color=DARK)

    # 5. 기대 효과 — 정성적
    qual_shape = _find_shape(slide, "직사각형 94")
    if qual_shape:
        quals = definition.get("expected_effects", {}).get("qualitative", [])
        _set_multiline_text(qual_shape, quals, font_size=Pt(9), bullet_char="•", color=DARK)

    # 6. 과제 추진 시 고려사항
    consider_shape = _find_shape(slide, "직사각형 91")
    if consider_shape:
        considerations = definition.get("considerations", [])
        _set_multiline_text(consider_shape, considerations, font_size=Pt(9), bullet_char="•", color=DARK)


# ── 슬라이드 3: 과제 설계서 ────────────────────────────────────────────────────

def _fill_design_slide(slide, design: dict, definition: dict | None = None):
    """과제 설계서 슬라이드에 데이터를 채웁니다."""
    project_number = ""
    project_title = design.get("project_title", "")
    created_date = ""
    author = ""

    if definition:
        project_number = definition.get("project_number", "")
        created_date = definition.get("created_date", "")
        author = definition.get("author", "")

    # 헤더 — 제목은 헤더 바에만 표시, 그룹 안 입력필드는 비움
    header_shape = _find_shape(slide, "직사각형 164")
    if header_shape:
        _set_text(header_shape, f"■ 과제번호/이름     {project_number}. {project_title}")

    # 그룹 165 안의 "과제번호/이름을 입력해주세요" 텍스트 제거 (겹침 방지)
    group165 = _find_shape(slide, "그룹 165")
    if group165 and hasattr(group165, 'shapes'):
        for sub in group165.shapes:
            if sub.has_text_frame:
                for para in sub.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""

    date_shape = _find_shape(slide, "직사각형 168")
    if date_shape:
        _set_text(date_shape, f"작성일: {created_date}   |   작성자: {author}")

    # 8. AI 기술 유형 — 체크박스 ✓ 표시
    _CHECKBOX_MAP = {
        # label text → checkbox shape name
        "정보 추출": "직사각형 34",
        "텍스트 생성": "직사각형 35",
        "대화형 인터페이스": "직사각형 36",
        "멀티모달 처리": "직사각형 57",
        "예측": "직사각형 134",
        "군집 · 분류": "직사각형 132",
        "최적화": "직사각형 147",
        "추천": "직사각형 145",
        "RPA": "직사각형 160",
        "OCR": "직사각형 158",
    }
    tech_types = design.get("ai_tech_info", {}).get("tech_types", [])
    tech_names_joined = " ".join(design.get("ai_tech_info", {}).get("tech_names", [])).upper()

    # tech_names에서 checked 자동 보정 (LLM이 누락한 체크 복구)
    _AUTO_CHECK_KEYWORDS = {
        "정보 추출": ["RAG", "검색", "추출", "RETRIEVAL"],
        "텍스트 생성": ["LLM", "GPT", "텍스트 생성"],
        "대화형 인터페이스": ["CHATBOT", "CONVERSATIONAL", "대화"],
        "멀티모달 처리": ["MULTIMODAL", "멀티모달"],
        "예측": ["예측", "PREDICTION", "FORECAST"],
        "군집 · 분류": ["ML MODEL", "ML 모델", "분류", "CLASSIFICATION", "CLUSTERING"],
        "최적화": ["최적화", "OPTIMIZATION"],
        "추천": ["추천", "RECOMMENDATION"],
        "RPA": ["RPA"],
        "OCR": ["OCR"],
    }
    checked_labels: set[str] = set()
    for tt in tech_types:
        for item in tt.get("checked", []):
            checked_labels.add(item)
        # 자동 보정: tech_names에 키워드가 있으면 체크 추가
        for sub in tt.get("sub_types", []):
            keywords = _AUTO_CHECK_KEYWORDS.get(sub, [])
            if any(kw.upper() in tech_names_joined for kw in keywords):
                checked_labels.add(sub)

    for label, cb_shape_name in _CHECKBOX_MAP.items():
        if label in checked_labels:
            cb_shape = _find_shape(slide, cb_shape_name)
            if cb_shape:
                cb_shape.fill.solid()
                cb_shape.fill.fore_color.rgb = PWC_RED
                # ✓ 텍스트 추가
                if cb_shape.has_text_frame:
                    tf = cb_shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = "✓"
                    run.font.size = Pt(8)
                    run.font.color.rgb = WHITE
                    run.font.bold = True

    # 8. AI 기술 유형 — 기술 이름 (공통 중복 제거, bullet 없이 쉼표 나열, 최대 8개)
    tech_names_shape = _find_shape(slide, "직사각형 37")
    if tech_names_shape:
        tech_names_raw = design.get("ai_tech_info", {}).get("tech_names", [])
        # 중복 제거
        seen = set()
        tech_deduped: list[str] = []
        for tn in tech_names_raw:
            if tn not in seen:
                seen.add(tn)
                tech_deduped.append(tn)
        if len(tech_deduped) > 8:
            tech_text = ", ".join(tech_deduped[:7]) + f" 외 {len(tech_deduped)-7}개"
        else:
            tech_text = ", ".join(tech_deduped)
        _set_text(tech_names_shape, tech_text, font_size=Pt(8), color=DARK)

    # 9. Input / Output (각각 최대 3줄)
    io_data = design.get("input_output", {})

    def _limit3(items: list[str]) -> list[str]:
        """최대 3개 항목만 표시."""
        return items[:3]

    io_88 = _find_shape(slide, "직사각형 88")
    if io_88:
        internals = _limit3(io_data.get("input_internal", []))
        _set_multiline_text(io_88, internals, font_size=Pt(8), bullet_char="•", color=DARK)

    io_89 = _find_shape(slide, "직사각형 89")
    if io_89:
        externals = _limit3(io_data.get("input_external", []))
        _set_multiline_text(io_89, externals, font_size=Pt(8), bullet_char="•", color=DARK)

    io_90 = _find_shape(slide, "직사각형 90")
    if io_90:
        outputs = _limit3(io_data.get("output", []))
        _set_multiline_text(io_90, outputs, font_size=Pt(8), bullet_char="•", color=DARK)


# ── 슬라이드 4: Agent 정의서 ──────────────────────────────────────────────────

def _fill_agent_slide(slide, agent: dict, design: dict, definition: dict | None = None):
    """Agent 정의서 슬라이드에 데이터를 채웁니다."""
    created_date = definition.get("created_date", "") if definition else ""
    author = definition.get("author", "") if definition else ""

    # 작성일/작성자
    date_shape = _find_shape(slide, "직사각형 168")
    if date_shape:
        _set_text(date_shape, f"작성일: {created_date}   |   작성자: {author}")

    # Agent 명
    name_shape = _find_shape(slide, "직사각형 5")
    if name_shape:
        _set_text(name_shape, agent.get("agent_name", ""))

    # Agent 유형 — Senior AI / Junior AI 표시
    # 타원 8 = Senior AI 옆 원, 타원 10 = Junior AI 옆 원
    # 선택된 쪽은 채움(●), 미선택은 비움(○)
    agent_type = agent.get("agent_type", "Junior AI")
    is_senior = agent_type == "Senior AI"

    senior_oval = _find_shape(slide, "타원 8")
    junior_oval = _find_shape(slide, "타원 10")

    if senior_oval:
        if is_senior:
            senior_oval.fill.solid()
            senior_oval.fill.fore_color.rgb = BLACK
        else:
            senior_oval.fill.background()
            senior_oval.line.color.rgb = GRAY
            senior_oval.line.width = Pt(1)

    if junior_oval:
        if not is_senior:
            junior_oval.fill.solid()
            junior_oval.fill.fore_color.rgb = BLACK
        else:
            junior_oval.fill.background()
            junior_oval.line.color.rgb = GRAY
            junior_oval.line.width = Pt(1)

    # Agent 역할
    roles = agent.get("roles", [])
    role_shape = _find_shape(slide, "직사각형 17")
    if role_shape:
        _set_multiline_text(role_shape, roles, font_size=Pt(9), bullet_char="•", color=DARK)

    # 처리 로직 — Input / 방법론 / Output
    input_data = agent.get("input_data", [])
    processing_steps = agent.get("processing_steps", [])
    output_data = agent.get("output_data", [])

    # ── 유사 항목 그룹핑 ──
    # "XXX 데이터 (세부1, 세부2)" → 괄호 앞 키워드를 기준으로 묶기
    # 괄호 없는 항목도 공통 접미사(데이터, 결과, 리포트 등)로 그룹핑
    import re as _re

    def _group_items(items: list[str], max_groups: int = 5) -> list[str]:
        """유사 항목을 큰 chunk로 묶어 '대분류 (세부1, 세부2)' 형태로 반환."""
        if len(items) <= max_groups:
            return [_shorten(s, 60) for s in items]

        # 1단계: 괄호가 있는 항목은 괄호 앞을 카테고리로 추출
        groups: dict[str, list[str]] = {}
        ungrouped: list[str] = []

        for item in items:
            s = str(item)
            # "AAA (BBB, CCC)" → category="AAA", detail="BBB, CCC"
            m = _re.match(r'^(.+?)\s*\((.+)\)\s*$', s)
            if m:
                cat = m.group(1).strip()
                detail = m.group(2).strip()
                groups.setdefault(cat, []).append(detail)
            else:
                # 공통 키워드로 그룹핑 시도 (데이터, 결과, 리포트, 목록 등)
                matched = False
                for keyword in ["데이터", "결과", "리포트", "보고서", "목록", "현황",
                                "분석", "전략", "지시", "브리핑", "요약"]:
                    if keyword in s:
                        groups.setdefault(keyword, []).append(s)
                        matched = True
                        break
                if not matched:
                    ungrouped.append(s)

        # 2단계: 그룹을 "카테고리 (세부1, 세부2)" 형태로 조합
        result: list[str] = []
        for cat, details in groups.items():
            if len(details) == 1:
                # 그룹에 1개면 괄호 안에 원래 있던 세부만
                combined = f"{cat} ({details[0][:40]})"
            else:
                # 여러 개 → 세부 나열
                short_details = [_shorten(d, 20) for d in details[:4]]
                suffix = f" 외 {len(details)-4}건" if len(details) > 4 else ""
                combined = f"{cat} ({', '.join(short_details)}{suffix})"
            result.append(_shorten(combined, 60))

        # 미분류 항목 추가
        for u in ungrouped[:max(1, max_groups - len(result))]:
            result.append(_shorten(u, 60))
        if len(ungrouped) > max(1, max_groups - len(groups)):
            remaining = len(ungrouped) - max(1, max_groups - len(groups))
            result.append(f"... 외 {remaining}건")

        return result[:max_groups]

    def _shorten(s: str, max_len: int = 60) -> str:
        """문자열을 최대 길이로 자르기 (넉넉하게)."""
        s = str(s).strip()
        return s[:max_len-1] + "…" if len(s) > max_len else s

    # 처리 로직 영역 좌표 (템플릿 Slide 3 기준)
    # Input 그룹(id=48):  L=1.8cm → 내용 L=3.5cm, T=12.3cm, W=5.8cm, H=5.0cm
    # 방법론 그룹(id=46): L=11.0cm → 내용 L=13.0cm, T=12.3cm, W=7.5cm, H=5.0cm
    # Output 그룹(id=47): L=22.3cm → 내용 L=24.2cm, T=12.3cm, W=7.5cm, H=5.0cm

    # 최대 줄수 제한 (칸 넘침 방지)
    MAX_LINES = 5

    _LOGIC_FONT = "Noto Sans KR"
    _LOGIC_SIZE = Pt(10)

    # Input — 그룹핑 + 최대 줄수
    if input_data:
        items = _group_items(input_data, max_groups=MAX_LINES)
        _add_multiline_textbox(
            slide, left=Cm(3.5), top=Cm(12.3),
            width=Cm(5.8), height=Cm(5.0),
            lines=items, font_size=_LOGIC_SIZE, bullet="•",
            line_spacing=1.15, font_name=_LOGIC_FONT,
        )

    # 방법론 — "번호. 이름 [기법] → 산출물" 최대 4줄
    if processing_steps:
        method_lines = []
        for idx, ps in enumerate(processing_steps[:4], 1):
            name = _shorten(str(ps.get("step_name", "")), 30)
            method = _shorten(str(ps.get("method", "")), 20)
            result_text = _shorten(str(ps.get("result", "")), 20)
            line = f"{idx}. {name}"
            if method:
                line += f" [{method}]"
            if result_text:
                line += f" → {result_text}"
            method_lines.append(line)
        if len(processing_steps) > 4:
            method_lines.append(f"... 외 {len(processing_steps)-4}단계")

        _add_multiline_textbox(
            slide, left=Cm(13.0), top=Cm(12.3),
            width=Cm(7.5), height=Cm(5.0),
            lines=method_lines,
            font_size=_LOGIC_SIZE, bullet="",
            line_spacing=1.15, font_name=_LOGIC_FONT,
        )

    # Output — 그룹핑 + 최대 줄수
    if output_data:
        items = _group_items(output_data, max_groups=MAX_LINES)
        _add_multiline_textbox(
            slide, left=Cm(24.2), top=Cm(12.3),
            width=Cm(7.5), height=Cm(5.0),
            lines=items, font_size=_LOGIC_SIZE, bullet="•",
            line_spacing=1.15, font_name=_LOGIC_FONT,
        )

    # Input → 방법론 → Output 사이 화살표는 템플릿에 이미 포함되어 있음


# ── 메인 함수 ─────────────────────────────────────────────────────────────────

def _remove_shapes_by_ids(slide, shape_ids: list[int]):
    """슬라이드에서 특정 shape_id의 도형을 제거합니다."""
    spTree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if shape.shape_id in shape_ids:
            sp = shape._element
            spTree.remove(sp)
            print(f"[ppt_exporter] 도형 제거: id={shape.shape_id} name={shape.name}")


def _insert_workflow_shapes(slide, workflow_cache: dict | None):
    """AI Service Flow를 PPT 도형으로 직접 그립니다 (수정 가능)."""
    if not workflow_cache:
        return
    try:
        # 기존 AI Service Flow 배경 도형 제거 (그룹 127)
        _remove_shapes_by_ids(slide, [128])

        from ppt_flow_drawer import draw_service_flow
        draw_service_flow(slide, workflow_cache)
        print("[ppt_exporter] AI Service Flow 도형 삽입 완료")
    except Exception as e:
        print(f"[ppt_exporter] AI Service Flow 도형 삽입 실패: {e}")


def export_ppt(
    definition: dict | None = None,
    design: dict | None = None,
    workflow: dict | None = None,
) -> io.BytesIO:
    """과제 정의서/설계서를 PPT로 내보냅니다. 표지/Thank You 슬라이드는 제거."""
    prs = Presentation(str(TEMPLATE_PATH))

    # 삭제할 슬라이드 인덱스 (역순으로 삭제)
    slides_to_remove = []

    # Slide 0: 표지 → 제거 대상
    slides_to_remove.append(0)

    # Slide 1: 과제 정의서
    if definition:
        _fill_definition_slide(prs.slides[1], definition)

    # Slide 2: 과제 설계서
    if design:
        _fill_design_slide(prs.slides[2], design, definition)
    # AI Service Flow 도형 그리기
    if workflow:
        _insert_workflow_shapes(prs.slides[2], workflow)

    # Slide 3+: Agent 정의서
    agents = design.get("agent_definitions", []) if design else []

    # Senior AI가 없으면 워크플로우 정보로 자동 생성하여 맨 앞에 추가
    has_senior = any(a.get("agent_type") == "Senior AI" for a in agents)
    if not has_senior and agents:
        process_name = workflow.get("process_name", "") if workflow else ""
        junior_names = [a.get("agent_name", "") for a in agents if a.get("agent_type") != "Senior AI"][:3]
        senior_def = {
            "agent_id": "senior-auto",
            "agent_name": f"{process_name} 오케스트레이터" if process_name else "Senior AI 오케스트레이터",
            "agent_type": "Senior AI",
            "roles": ["Junior AI Agent 실행 오케스트레이션", "전체 워크플로우 조율 및 결과 품질 관리"],
            "input_data": [f"Junior AI 실행 결과 ({', '.join(junior_names)})" if junior_names else "Junior AI 실행 결과", "업무 요청 데이터"],
            "processing_steps": [
                {"step_name": "워크플로우 분석", "method": "LLM", "result": "실행 계획"},
                {"step_name": "Junior AI 지시·조율", "method": "오케스트레이션", "result": "실행 결과"},
                {"step_name": "결과 품질 검증", "method": "LLM", "result": "검증 보고"},
            ],
            "output_data": ["Junior AI 실행 지시", "최종 결과 품질 검증 보고"],
        }
        agents = [senior_def] + agents

    if agents:
        # 1) 먼저 필요한 슬라이드 수만큼 복제 (빈 상태에서 복제)
        agent_slides = [prs.slides[3]]  # 첫 번째는 원본
        for _ in agents[1:]:
            agent_slides.append(_duplicate_slide(prs, 3))

        # 2) 그 다음 각 슬라이드에 개별 Agent 데이터 채우기
        for slide, agent in zip(agent_slides, agents):
            _remove_shapes_by_ids(slide, [133])  # 기존 미니맵 배경 제거
            _fill_agent_slide(slide, agent, design, definition)
            if workflow:
                try:
                    from ppt_flow_drawer import draw_minimap
                    draw_minimap(slide, workflow, highlight_agent_id=agent.get("agent_id", ""))
                except Exception as e:
                    print(f"[ppt_exporter] 미니맵 실패: {e}")

    # Thank You 슬라이드 (마지막) → 제거 대상
    thank_you_idx = 4 if not agents else 4  # 복제 전 기준
    slides_to_remove.append(thank_you_idx)

    # 슬라이드 제거 (역순)
    for idx in sorted(slides_to_remove, reverse=True):
        if idx < len(prs.slides):
            rId = prs.slides._sldIdLst[idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', '')
            sldId = prs.slides._sldIdLst[idx]
            prs.slides._sldIdLst.remove(sldId)

    # BytesIO로 저장
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def _move_slide_to_end(prs: Presentation, slide_index: int):
    """슬라이드를 맨 뒤로 이동합니다 (XML 기반)."""
    slides_el = prs.slides._sldIdLst
    slide_ids = list(slides_el)
    if slide_index < len(slide_ids):
        target = slide_ids[slide_index]
        slides_el.remove(target)
        slides_el.append(target)
