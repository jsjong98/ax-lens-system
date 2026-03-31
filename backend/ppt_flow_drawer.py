"""
ppt_flow_drawer.py вҖ” PPTм—җ AI Service FlowлҘј лҸ„нҳ•мңјлЎң м§Ғм ‘ к·ёлҰ¬кё°

python-pptx лҸ„нҳ•(мӮ¬к°Ғнҳ•)л§Ң мӮ¬мҡ©н•ҳм—¬ мҲҳм • к°ҖлҠҘн•ң мҠӨмң”л Ҳмқёмқ„ к·ёлҰҪлӢҲлӢӨ.
вҖ» python-pptx add_connector лҢҖмӢ  OOXML cxnSp XMLмқ„ м§Ғм ‘ мғқм„ұн•ҳм—¬ PPT ліөкө¬ мҳӨлҘҳлҘј л°©м§Җн•©лӢҲлӢӨ.
- draw_service_flow: м „мІҙ AI Service Flow (кіјм ң м„Өкі„м„ңмҡ©)
- draw_minimap: м¶•мҶҢ лҜёлӢҲл§ө (Agent м •мқҳм„ңмҡ©, нҠ№м • Agent к°•мЎ°)
"""
from __future__ import annotations

from pptx.util import Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# в”Җв”Җ мғүмғҒ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ

RED = RGBColor(0x8B, 0x1A, 0x1A)
RED_BG = RGBColor(0xF8, 0xE0, 0xE0)
GOLD = RGBColor(0xAA, 0x8E, 0x2A)
GRAY = RGBColor(0x88, 0x87, 0x80)
GRAY_ARROW = RGBColor(0x9E, 0x9E, 0x9E)
LIGHT_GRAY = RGBColor(0xD3, 0xD1, 0xC7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x2C, 0x2C, 0x2A)
YELLOW_BG = RGBColor(0xFE, 0xFA, 0xF0)
GRAY_BG = RGBColor(0xF5, 0xF5, 0xF3)
INPUT_BG = RGBColor(0xF5, 0xF4, 0xF1)

# Agentлі„ кө¬л¶„ мғүмғҒ вҖ” нҢҢлһҖ кі„м—ҙ мң м§Җ, лӘ…лҸ„В·мұ„лҸ„ м°ЁмқҙлЎң кө¬л¶„
_AGENT_PALETTE = [
    RGBColor(0x1A, 0x3C, 0x6E),   # 1  м§„лӮЁмғү (к°ҖмһҘ м–ҙл‘җмӣҖ)
    RGBColor(0x2E, 0x75, 0xB6),   # 2  мӨ‘к°„ нҢҢлһҖ
    RGBColor(0x00, 0x82, 0x7F),   # 3  нӢё (мҙҲлЎқ кё°мҡҙ)
    RGBColor(0x5B, 0x9B, 0xD5),   # 4  л°қмқҖ н•ҳлҠҳ
    RGBColor(0x4B, 0x00, 0x82),   # 5  мқёл””кі  (ліҙлқј кё°мҡҙ)
    RGBColor(0x00, 0xA6, 0xA0),   # 6  л°қмқҖ мІӯлЎқ
    RGBColor(0x41, 0x72, 0xC4),   # 7  мҪ”л°ңнҠё
    RGBColor(0x7B, 0x68, 0xC4),   # 8  нҚјн”Ңлё”лЈЁ
    RGBColor(0x00, 0x6E, 0x90),   # 9  нҺҳнҠёлЎӨ
    RGBColor(0x87, 0xCE, 0xEB),   # 10 мҠӨм№ҙмқҙ (к°ҖмһҘ л°қмқҢ)
]


def _agent_color(idx: int) -> RGBColor:
    return _AGENT_PALETTE[idx % len(_AGENT_PALETTE)]


# в”Җв”Җ кё°ліё лҸ„нҳ• н—¬нҚј в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ

def _add_rect(slide, left, top, width, height, fill=None, border_color=None,
              border_width=Pt(1), border_dash=None, text="", font_size=Pt(8),
              font_color=BLACK, bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or WHITE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
        if border_dash:
            shape.line.dash_style = border_dash
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = bold
    return shape


# в”Җв”Җ л Ҳмқё к°„ м—°кІ° мғүмғҒ (м°ёмЎ° н…ңн”ҢлҰҝ кё°мӨҖ) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
CONN_BLUE = RGBColor(0x4D, 0xAC, 0xF1)    # InputвҶ’Senior, SeniorвҶ’Junior
CONN_GOLD = RGBColor(0xB4, 0x8E, 0x04)    # JuniorвҶ’HR
CONN_RED = RGBColor(0x8B, 0x1A, 0x1A)     # SeniorвҶ’HR (к°җлҸ…)
CONN_GRAY = RGBColor(0x99, 0x99, 0x99)    # JuniorвҶ’Senior (н”јл“ңл°ұ)
CONN_WIDTH = 12700                          # 1pt in EMU


def _make_cxnSp(slide, x1, y1, x2, y2, color=CONN_BLUE, width=CONN_WIDTH,
                 head_end=False, tail_end=True, bent=False):
    """м°ёмЎ° н…ңн”ҢлҰҝкіј лҸҷмқјн•ң OOXML cxnSp м»Өл„Ҙн„°лҘј мғқм„ұн•©лӢҲлӢӨ.
    bent=Trueмқҙл©ҙ кәҫмқё м»Өл„Ҙн„°(bentConnector3)лҘј мӮ¬мҡ©н•©лӢҲлӢӨ."""
    from lxml import etree
    from pptx.oxml.ns import qn

    flip_h = x2 < x1
    flip_v = y2 < y1
    off_x = min(x1, x2)
    off_y = min(y1, y2)
    cx = abs(x2 - x1)
    cy = abs(y2 - y1)

    sp_tree = slide.shapes._spTree
    max_id = max(
        (int(el.get('id', '0'))
         for el in sp_tree.iter()
         if el.get('id', '').isdigit()),
        default=0,
    )
    new_id = max_id + 1
    color_hex = str(color)

    cxn = etree.SubElement(sp_tree, qn('p:cxnSp'))

    # nvCxnSpPr
    nv = etree.SubElement(cxn, qn('p:nvCxnSpPr'))
    cNvPr = etree.SubElement(nv, qn('p:cNvPr'))
    cNvPr.set('id', str(new_id))
    cNvPr.set('name', f'Connector {new_id}')
    cNvCxnSpPr = etree.SubElement(nv, qn('p:cNvCxnSpPr'))
    etree.SubElement(cNvCxnSpPr, qn('a:cxnSpLocks'))
    etree.SubElement(nv, qn('p:nvPr'))

    # spPr
    spPr = etree.SubElement(cxn, qn('p:spPr'))

    xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    if flip_h:
        xfrm.set('flipH', '1')
    if flip_v:
        xfrm.set('flipV', '1')

    off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(int(off_x)))
    off.set('y', str(int(off_y)))
    ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(int(cx)))
    ext.set('cy', str(int(cy)))

    prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
    if bent:
        prstGeom.set('prst', 'bentConnector3')
        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        # кәҫмһ„ м§Җм җ: 50000 = мӨ‘к°„ (0~100000 лІ”мң„)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj1')
        gd.set('fmla', 'val 50000')
    else:
        prstGeom.set('prst', 'straightConnector1')
        etree.SubElement(prstGeom, qn('a:avLst'))

    ln = etree.SubElement(spPr, qn('a:ln'))
    ln.set('w', str(int(width)))
    ln.set('cap', 'sq')

    solidFill = etree.SubElement(ln, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', color_hex)

    if tail_end:
        te = etree.SubElement(ln, qn('a:tailEnd'))
        te.set('type', 'triangle')
    if head_end:
        he = etree.SubElement(ln, qn('a:headEnd'))
        he.set('type', 'triangle')

    # p:style (м°ёмЎ° н…ңн”ҢлҰҝ лҸҷмқј)
    style = etree.SubElement(cxn, qn('p:style'))
    for ref_tag, ref_idx, clr_val in [
        ('a:lnRef', '1', 'accent1'),
        ('a:fillRef', '0', 'accent1'),
        ('a:effectRef', '0', 'dk1'),
    ]:
        ref = etree.SubElement(style, qn(ref_tag))
        ref.set('idx', ref_idx)
        sc = etree.SubElement(ref, qn('a:schemeClr'))
        sc.set('val', clr_val)
    fontRef = etree.SubElement(style, qn('a:fontRef'))
    fontRef.set('idx', 'minor')
    sc = etree.SubElement(fontRef, qn('a:schemeClr'))
    sc.set('val', 'lt1')


def _arrow_v(slide, x, y1, y2, color=CONN_BLUE, width=CONN_WIDTH):
    """м„ёлЎң нҷ”мӮҙн‘ң (y1вҶ’y2 л°©н–Ҙ)."""
    _make_cxnSp(slide, x, y1, x, y2, color=color, width=width, tail_end=True)


def _draw_vline(slide, x, y1, y2, color=CONN_BLUE, width=CONN_WIDTH):
    """м„ёлЎңм„  (нҷ”мӮҙн‘ң лЁёлҰ¬ м—ҶмқҢ)."""
    _make_cxnSp(slide, x, y1, x, y2, color=color, width=width, tail_end=False)


# в•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җ
# лҜёлӢҲл§ө к·ёлҰ¬кё° (Agent м •мқҳм„ңмҡ©)
# н…ңн”ҢлҰҝ к·ёлЈ№132: L=24.4cm, T=3.7cm, W=7.8cm, H=5.7cm
# в•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җ

def draw_minimap(slide, workflow: dict, highlight_agent_id: str = "",
                 left=Cm(24.4), top=Cm(3.7), total_width=Cm(7.8), total_height=Cm(5.7)):
    """AI Service FlowмҷҖ лҸҷмқјн•ң кө¬мЎ°мқҳ м¶•мҶҢ лҜёлӢҲл§ө.
    - Input л°•мҠӨ: мӢӨм ң Input мҲҳл§ҢнҒј н‘ңмӢң, Agentлі„ мғүмғҒ
    - кәҫмқём„ : Input вҶ’ Junior AI м§Ғм ‘ м—°кІ° (bent connector)
    - Senior AI л°”к°Җ нҷ”мӮҙн‘ңлҘј лҚ®мқҢ
    """
    agents = workflow.get("agents", [])
    if not agents:
        return

    agent_count = len(agents)
    label_w = Cm(1.0)
    content_w = total_width - label_w
    agent_col_w = content_w / max(agent_count, 1)
    content_left = left + label_w

    # л Ҳмқё лҶ’мқҙ л°°л¶„
    input_h  = Cm(0.7)
    gap1     = Cm(0.25)
    senior_h = Cm(0.7)
    gap2     = Cm(0.3)
    junior_h = Cm(2.2)
    gap3     = Cm(0.25)
    hr_h     = Cm(0.7)

    input_top  = top + Cm(0.15)
    senior_top = input_top + input_h + gap1
    junior_top = senior_top + senior_h + gap2
    hr_top     = junior_top + junior_h + gap3

    # в”Җв”Җ Agentлі„ Input мҲҳм§‘ (AI Service FlowмҷҖ лҸҷмқј лЎңм§Ғ) в”Җв”Җ
    agent_inputs: list[list[str]] = []
    input_owner: dict[str, int] = {}
    for ai, a in enumerate(agents):
        inps: list[str] = []
        for t in a.get("assigned_tasks", []):
            for inp in t.get("input_data", []):
                if inp not in inps:
                    inps.append(inp)
                if inp not in input_owner:
                    input_owner[inp] = ai
        agent_inputs.append(inps)

    all_inps = list(dict.fromkeys(i for ai in agent_inputs for i in ai))[:10]
    inp_w = content_w / max(len(all_inps), 1)
    inp_cx: dict[str, int] = {}

    # в”Җв”Җ Input н–ү в”Җв”Җ
    _add_rect(slide, left, input_top, label_w, input_h,
              fill=WHITE, border_color=LIGHT_GRAY,
              text="Input", font_size=Pt(4), font_color=GRAY)
    for i, inp in enumerate(all_inps):
        bx = content_left + i * inp_w + Cm(0.02)
        bw = inp_w - Cm(0.04)
        oi = input_owner.get(inp, 0)
        c = _agent_color(oi)
        _add_rect(slide, bx, input_top + Cm(0.06), bw, input_h - Cm(0.12),
                  fill=WHITE, border_color=c, border_width=Pt(0.5))
        inp_cx[inp] = bx + bw // 2

    # Agentлі„ м»¬лҹј мӨ‘м•ҷ X
    agent_cxs: list[int] = []
    for i in range(agent_count):
        cl = content_left + i * agent_col_w
        agent_cxs.append(cl + agent_col_w // 2)

    # в”Җв”Җ InputвҶ’Junior кәҫмқё нҷ”мӮҙн‘ң (лЁјм Җ к·ёл Өм„ң Senior л°” л’ӨлЎң) в”Җв”Җ
    mini_w = Emu(6350)  # 0.5pt вҖ” лҜёлӢҲл§өмҡ© м–ҮмқҖ м„ 
    for i, agent in enumerate(agents):
        cx = agent_cxs[i]
        c = _agent_color(i)
        inps = agent_inputs[i]
        for inp in inps:
            sx = inp_cx.get(inp)
            if sx is not None:
                _make_cxnSp(slide, sx, input_top + input_h, cx, junior_top,
                            color=c, width=mini_w, bent=True)

    # в”Җв”Җ Senior AI н–ү (нҷ”мӮҙн‘ң мң„м—җ лҚ®мқҢ) в”Җв”Җ
    _add_rect(slide, left, senior_top, label_w, senior_h,
              fill=WHITE, border_color=LIGHT_GRAY,
              text="Senior\nAI", font_size=Pt(4), font_color=RED)
    _add_rect(slide, content_left, senior_top + Cm(0.04),
              content_w, senior_h - Cm(0.08),
              fill=RED_BG, border_color=RED, border_width=Pt(1.5))

    # в”Җв”Җ SeniorвҶ”Junior нҷ”мӮҙн‘ң в”Җв”Җ
    for i in range(agent_count):
        cx = agent_cxs[i]
        _arrow_v(slide, cx - Cm(0.06), senior_top + senior_h, junior_top, CONN_RED, mini_w)
        _arrow_v(slide, cx + Cm(0.06), junior_top, senior_top + senior_h, CONN_GRAY, mini_w)

    # в”Җв”Җ Junior AI н–ү в”Җв”Җ
    _add_rect(slide, left, junior_top, label_w, junior_h,
              fill=WHITE, border_color=LIGHT_GRAY,
              text="Junior\nAI", font_size=Pt(4), font_color=GOLD)

    for i, agent in enumerate(agents):
        bx = content_left + i * agent_col_w + Cm(0.05)
        bw = agent_col_w - Cm(0.1)
        is_hl = agent.get("agent_id") == highlight_agent_id

        _add_rect(slide, bx, junior_top + Cm(0.05), bw, junior_h - Cm(0.1),
                  fill=YELLOW_BG if is_hl else WHITE,
                  border_color=GOLD,
                  border_width=Pt(2) if is_hl else Pt(0.4))

        # лІҲнҳё л°°м§Җ
        _add_rect(slide, bx + Cm(0.04), junior_top + Cm(0.1), Cm(0.3), Cm(0.3),
                  fill=GOLD if is_hl else LIGHT_GRAY,
                  text=str(i + 1), font_size=Pt(4),
                  font_color=WHITE if is_hl else BLACK, bold=True)

        # Task л°•мҠӨ
        tasks = agent.get("assigned_tasks", [])
        for j in range(min(len(tasks), 3)):
            ty = junior_top + Cm(0.5) + j * Cm(0.4)
            _add_rect(slide, bx + Cm(0.06), ty,
                      bw - Cm(0.12), Cm(0.3),
                      fill=INPUT_BG, border_color=GOLD, border_width=Pt(0.2),
                      border_dash=7)

    # в”Җв”Җ HR н–ү в”Җв”Җ
    _add_rect(slide, left, hr_top, label_w, hr_h,
              fill=WHITE, border_color=LIGHT_GRAY,
              text="HR\nлӢҙлӢ№мһҗ", font_size=Pt(4), font_color=BLACK)
    for i, agent in enumerate(agents):
        has_hr = any(t.get("human_role", "") for t in agent.get("assigned_tasks", []))
        if has_hr:
            bx = content_left + i * agent_col_w + Cm(0.05)
            bw = agent_col_w - Cm(0.1)
            _add_rect(slide, bx, hr_top + Cm(0.05), bw, hr_h - Cm(0.1),
                      fill=GRAY_BG, border_color=GRAY_ARROW, border_width=Pt(0.3))

    # в”Җв”Җ JuniorвҶ’HR нҷ”мӮҙн‘ң (кёҲмғү, HR мһҲлҠ” кІҪмҡ°л§Ң) в”Җв”Җ
    for i, agent in enumerate(agents):
        cx = agent_cxs[i]
        has_hr = any(t.get("human_role", "") for t in agent.get("assigned_tasks", []))
        if has_hr:
            _arrow_v(slide, cx, junior_top + junior_h, hr_top, CONN_GOLD, mini_w)



# в•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җ
# м „мІҙ AI Service Flow к·ёлҰ¬кё° (кіјм ң м„Өкі„м„ңмҡ©)
# н…ңн”ҢлҰҝ мҳҒм—ӯ: L=1.0cm, T=3.7cm, W=18.3cm, H=13.5cm
# в•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җв•җ

def draw_service_flow(slide, workflow: dict,
                      left=Cm(1.0), top=Cm(3.7), total_width=Cm(18.3), total_height=Cm(13.5)):
    agents = workflow.get("agents", [])
    if not agents:
        return

    agent_count = len(agents)
    label_w = Cm(1.2)
    content_w = total_width - label_w
    agent_col_w = content_w / max(agent_count, 1)
    content_left = left + label_w

    # л Ҳмқё лҶ’мқҙ (м „мІҙ 13.5cmм—җ л§һм¶Ө)
    input_h  = Cm(1.2)
    gap_is   = Cm(0.5)
    senior_h = Cm(1.2)
    gap_sj   = Cm(0.6)
    junior_h = Cm(7.5)
    gap_jh   = Cm(0.5)
    hr_h     = Cm(1.2)
    # н•©кі„: 1.2+0.5+1.2+0.6+7.5+0.5+1.2 = 12.7cm < 13.5cm вң“

    r = [top]
    r.append(r[0] + input_h + gap_is)
    r.append(r[1] + senior_h + gap_sj)
    r.append(r[2] + junior_h + gap_jh)

    # в”Җв”Җ л Ҳмқё л°°кІҪ в”Җв”Җ
    _add_rect(slide, content_left, r[2], content_w, junior_h,
              fill=YELLOW_BG, border_color=None)

    # в•җв•җ 1. Input н–ү в•җв•җ
    _add_rect(slide, left, r[0], label_w, input_h,
              fill=WHITE, border_color=LIGHT_GRAY,
              text="Input", font_size=Pt(6), font_color=GRAY)

    # Agentлі„ Input мҲҳм§‘
    agent_inputs: list[list[str]] = []
    input_owner: dict[str, int] = {}
    for ai, a in enumerate(agents):
        inps: list[str] = []
        for t in a.get("assigned_tasks", []):
            for inp in t.get("input_data", []):
                if inp not in inps:
                    inps.append(inp)
                if inp not in input_owner:
                    input_owner[inp] = ai
        agent_inputs.append(inps)

    all_inps = list(dict.fromkeys(i for ai in agent_inputs for i in ai))[:10]
    inp_w = content_w / max(len(all_inps), 1)
    inp_cx: dict[str, int] = {}

    for i, inp in enumerate(all_inps):
        bx = content_left + i * inp_w + Cm(0.08)
        bw = inp_w - Cm(0.16)
        oi = input_owner.get(inp, 0)
        c = _agent_color(oi)
        _add_rect(slide, bx, r[0] + Cm(0.1), bw, input_h - Cm(0.2),
                  fill=WHITE, border_color=c, border_width=Pt(1),
                  text=inp[:15], font_size=Pt(5), font_color=BLACK)
        inp_cx[inp] = bx + bw // 2

    # Agentлі„ м»¬лҹј мӨ‘м•ҷ X мўҢн‘ң (нҷ”мӮҙн‘ңм—җ н•„мҡ”)
    agent_cxs: list[int] = []
    for i in range(agent_count):
        cl = content_left + i * agent_col_w
        agent_cxs.append(cl + agent_col_w // 2)

    # в•җв•җ 2. InputвҶ’Junior кәҫмқё нҷ”мӮҙн‘ң (лЁјм Җ к·ёл Өм„ң Senior AI л°” л’ӨлЎң) в•җв•җ
    for i, agent in enumerate(agents):
        cx = agent_cxs[i]
        c = _agent_color(i)
        inps = agent_inputs[i]
        for inp in inps:
            sx = inp_cx.get(inp)
            if sx is not None:
                _make_cxnSp(slide, sx, r[0] + input_h, cx, r[2],
                            color=c, bent=True)

    # в•җв•җ 3. Senior AI н–ү (нҷ”мӮҙн‘ң мң„м—җ лҚ®мқҢ) в•җв•җ
    _add_rect(slide, left, r[1], label_w, senior_h,
              fill=WHITE, border_color=LIGHT_GRAY,
              text="Senior\nAI", font_size=Pt(6), font_color=RED)

    pname = workflow.get("process_name", "")
    _add_rect(slide, content_left + Cm(0.15), r[1] + Cm(0.1),
              content_w - Cm(0.3), senior_h - Cm(0.2),
              fill=RED_BG, border_color=RED, border_width=Pt(1.2),
              text=f"{pname} мҳӨмјҖмҠӨнҠёл Ҳмқҙн„°", font_size=Pt(6), font_color=RED, bold=True)

    # в•җв•җ 4. SeniorвҶ”Junior нҷ”мӮҙн‘ң (Senior л°” м•„лһҳм—җ) в•җв•җ
    for i in range(agent_count):
        cx = agent_cxs[i]
        # Senior AI вҶ’ Junior AI (л№Ёк°„)
        _arrow_v(slide, cx - Cm(0.12), r[1] + senior_h, r[2], CONN_RED)
        # Junior AI вҶ’ Senior AI н”јл“ңл°ұ (нҡҢмғү)
        _arrow_v(slide, cx + Cm(0.12), r[2], r[1] + senior_h, CONN_GRAY)

    # в•җв•җ 5. Junior AI н–ү в•җв•җ
    _add_rect(slide, left, r[2], label_w, junior_h,
              fill=WHITE, border_color=LIGHT_GRAY,
              text="Junior\nAI", font_size=Pt(6), font_color=GOLD)

    for i, agent in enumerate(agents):
        cl = content_left + i * agent_col_w
        bx = cl + Cm(0.1)
        bw = agent_col_w - Cm(0.2)

        _add_rect(slide, bx, r[2] + Cm(0.2), bw, junior_h - Cm(0.3),
                  fill=YELLOW_BG, border_color=GOLD, border_width=Pt(1))

        # лІҲнҳё
        _add_rect(slide, bx + Cm(0.06), r[2] + Cm(0.28), Cm(0.4), Cm(0.4),
                  fill=GOLD, text=str(i + 1), font_size=Pt(6), font_color=WHITE, bold=True)

        # Task л°•мҠӨ
        tasks = agent.get("assigned_tasks", [])
        task_h = Cm(0.5)
        task_gap = Cm(0.2)
        max_t = min(len(tasks), int((junior_h - Cm(1.2)) / (task_h + task_gap)))

        for j in range(max_t):
            t = tasks[j]
            ty = r[2] + Cm(0.85) + j * (task_h + task_gap)
            _add_rect(slide, bx + Cm(0.1), ty, bw - Cm(0.2), task_h,
                      fill=INPUT_BG, border_color=GOLD, border_width=Pt(0.5),
                      border_dash=7,
                      text=t.get("task_name", "")[:18], font_size=Pt(4.5), font_color=BLACK)
            if j > 0:
                _arrow_v(slide, bx + bw / 2, ty - task_gap, ty, CONN_GOLD, Emu(9525))

    # в•җв•җ 6. HR н–ү в•җв•җ
    _add_rect(slide, left, r[3], label_w, hr_h,
              fill=WHITE, border_color=LIGHT_GRAY,
              text="HR\nлӢҙлӢ№мһҗ", font_size=Pt(6), font_color=BLACK)

    for i, agent in enumerate(agents):
        cl = content_left + i * agent_col_w
        human_tasks = [t for t in agent.get("assigned_tasks", []) if t.get("human_role", "")]
        if human_tasks:
            _add_rect(slide, cl + Cm(0.1), r[3] + Cm(0.1),
                      agent_col_w - Cm(0.2), hr_h - Cm(0.2),
                      fill=WHITE, border_color=GRAY_ARROW, border_width=Pt(0.7),
                      text=human_tasks[0].get("human_role", "кІҖнҶ ")[:18],
                      font_size=Pt(4.5), font_color=BLACK)

    # в•җв•җ 7. JuniorвҶ’HR нҷ”мӮҙн‘ң (кёҲмғү, HR мһҲлҠ” кІҪмҡ°л§Ң) в•җв•җ
    for i, agent in enumerate(agents):
        cx = agent_cxs[i]
        has_hr = any(t.get("human_role", "") for t in agent.get("assigned_tasks", []))
        if has_hr:
            _arrow_v(slide, cx, r[2] + junior_h, r[3], CONN_GOLD)
